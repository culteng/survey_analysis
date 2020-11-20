# use python-pptx to generate ppt report

import glob
import sys
from datetime import datetime
import pandas as pd
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
prs = Presentation('./template/ppt_template_v01.pptx')

# project imports
import analysis


# functions for pptx
def scale(image, max_size, method=Image.ANTIALIAS):
    # add whitespace to images to avoid cropping for fit in placeholders
    image.thumbnail(max_size, method)
    offset = (int((max_size[0] - image.size[0]) / 2), int((max_size[1] - image.size[1]) / 2))
    back = Image.new("RGB", max_size, "white")
    back.paste(image, offset)
    return back


def func(pct, allvals):
    # back-calculates # of 2ks from pct given by pie 
    absolute = int(round(pct/100.*np.sum(allvals)))
    return absolute


def font_plainerizer(table):
    # removes powerpoint default table formatting
    for c in range(2):
        for r in range(2):
            cell = table.cell(r-1, c-1)
            cell.fill.solid()
            cell.fill.fore_color.rgb=RGBColor(255,255,255)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(28,41,91)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.size=Pt(14)


def df_to_ppt(slide, df, target_ph, col_widths, header = True):
    # cleans & loads pd.df to pptx placeholder
    placeholder_i = slide.placeholders[target_ph]
    ''' can't sort....
    fill = placeholder_i.fill
    fill = None'''
    graphic_frame = placeholder_i.insert_table(rows=(df.shape[0] + 1), cols=df.shape[1])
    table = graphic_frame.table
    # font_plainerizer(table)
    for idx, wid in enumerate(col_widths):
        table.columns[idx].width = Inches(wid)
    for rows in range(df.shape[0] + 1):
        table.rows[rows].height = Inches(.32)
        for cols in range(df.shape[1]):
            if rows == 0 and header:
                table.cell(rows, cols).text = list(df.columns.values)[cols]
                table.cell(rows, cols).text_frame.paragraphs[0].font.size=Pt(15)
            else:
                dat = df.iloc[rows-1, cols]
                table.cell(rows, cols).text = str(int(round(dat)) if isinstance(dat, float) else dat)
                table.cell(rows, cols).text_frame.paragraphs[0].font.size=Pt(8)
                if target_ph == 21: table.cell(rows, cols).text_frame.paragraphs[0].font.size=Pt(8)


def slide_title():
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    placeholder = slide.placeholders[22]
    txfr = placeholder.text_frame
    txfr.paragraphs[0].alignment = PP_ALIGN.CENTER
    placeholder.text = school_name


def main(args):
    print("report.py main starting")
    filenames = glob.glob("data/cleaned/*.csv")
    df = analysis.populate_df(filenames)
    
    # aggregate stats for likert questions
    df_likerts = df.loc[:, df.columns.str.startswith('scl')]
    df_counts = pd.get_dummies(df_likerts.stack()).groupby(level=1).sum()
    df_counts.columns = ['Strongly Disagree', 'Disagree', 'Neutral', 'Agree', 'Strongly Agree']
    df_likert_stats = df_likerts.join(df_likerts.mean(1).rename('mean')).join(df_likerts.std(1).rename('std'))
    
    # cover page 
    school_name = 'CESA 1'
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # insert logo
    placeholder = slide.placeholders[13]
    placeholder.insert_picture('./imgs/cesa-logo.png')
    # insert school name
    placeholder = slide.placeholders[2]
    placeholder.text = 'Council on Educational Standards and Accountability'
    # fetch contact details
    school_contact = {
        'name': 'Yon Yonson',
        'title': 'President',
        'phone': '(123) 456-7890',
        'email': 'yon@cesa.edu'}
    # insert contact details
    placeholder = slide.placeholders[16]
    placeholder.text = '{0}\n{1}'.format(school_contact['name'], school_contact['title'])
    placeholder = slide.placeholders[17]
    placeholder.text = '{0}\n{1}'.format(school_contact['phone'], school_contact['email'])
    
    # exec summary page
    slide = prs.slides.add_slide(prs.slide_layouts[1])                
    df_class_counts = df.groupby(('i_am_a')).size().to_frame('count')
    df_to_ppt(slide, df_class_counts, 13, [1], header=False)
    # insert key findings
    placeholder = slide.placeholders[14]
    # fetch findings
    findings = ['The data revealed that Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud exercitation',
        'analysis of ...Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt ',
        'while the intial view of ....Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam'
        ]
    placeholder.text = '\n'.join(findings)
    
    
    # analysis by section - convert to function over each group
    # df_to_ppt(df_likert_stats.iloc[:5][['mean','std']])
    
    prs.save('./output/Strickland Report - {}-{}.pptx'.format(school_name, datetime.now().strftime('%Y%m%d')))
    print('report.py done')


if __name__ == '__main__':
    main(sys.argv[1:])


''' utils/samples
# get placeholders from current slide
for shape in slide.placeholders:
    print(shape.placeholder_format.idx, shape.name)
'''

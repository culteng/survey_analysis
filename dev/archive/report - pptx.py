# use python-pptx to generate ppt report

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
prs = Presentation('./report_temp_div3.pptx')


# functions for pptx
def scale(image, max_size, method=Image.ANTIALIAS):
    # add whitespace to images to avoid cropping for fit in placeholders
    image.thumbnail(max_size, method)
    offset = (int((max_size[0] - image.size[0]) / 2), int((max_size[1] - image.size[1]) / 2))
    back = Image.new("RGB", max_size, "white")
    back.paste(image, offset)
    return back


def func(pct, allvals):
    # back-calculates # of 2ks from pct given by pie 
    absolute = int(round(pct/100.*np.sum(allvals)))
    return absolute


def df_to_ppt(df, target_ph, col_widths, header = True):
    # cleans & loads pd.df to pptx placeholder
    placeholder_i = slide.placeholders[target_ph]
    graphic_frame = placeholder_i.insert_table(rows=(df.shape[0] + 1), cols=df.shape[1])
    table = graphic_frame.table
    for idx, wid in enumerate(col_widths):
        table.columns[idx].width = Inches(wid)
    for rows in range(df.shape[0] + 1):
        table.rows[rows].height = Inches(.32)
        for cols in range(df.shape[1]):
            if rows == 0 and header:
                table.cell(rows, cols).text = list(df.columns.values)[cols]
                table.cell(rows, cols).text_frame.paragraphs[0].font.size=Pt(15)
            else:
                dat = df.iloc[rows-1, cols]
                table.cell(rows, cols).text = str(int(round(dat)) if isinstance(dat, float) else dat)
                table.cell(rows, cols).text_frame.paragraphs[0].font.size=Pt(11)
                if target_ph == 21: table.cell(rows, cols).text_frame.paragraphs[0].font.size=Pt(8)


def font_plainerizer():
    # removes powerpoint default table formatting
    for c in range(2):
        for r in range(2):
            cell = table.cell(r-1, c-1)
            cell.fill.solid()
            cell.fill.fore_color.rgb=RGBColor(255,255,255)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(28,41,91)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.size=Pt(14)


def slide_title():
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    placeholder = slide.placeholders[22]
    txfr = placeholder.text_frame
    txfr.paragraphs[0].alignment = PP_ALIGN.CENTER
    placeholder.text = school_name


def main(args):
    print("report.py main")


if __name__ == '__main__':
    main(sys.argv[1:])

